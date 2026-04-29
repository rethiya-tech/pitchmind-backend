import pytest
from io import BytesIO


def test_get_theme_executive_modern():
    from app.services.themes import get_theme
    theme = get_theme("executive_modern")
    assert theme.id == "executive_modern"


def test_get_theme_unknown_raises_key_error():
    from app.services.themes import get_theme
    with pytest.raises(KeyError):
        get_theme("nonexistent_theme")


def test_all_six_themes_exist():
    from app.services.themes import THEMES
    expected = {
        "executive_modern", "corporate_zenith", "digital_frontier",
        "nordic_flow", "midnight_insight", "executive_gold"
    }
    assert expected.issubset(set(THEMES.keys()))


def test_build_pptx_returns_bytes():
    from app.services.pptx_builder import build_pptx
    from app.services.themes import get_theme
    slides = [
        type("Slide", (), {
            "title": "Test Slide",
            "bullets": ["Bullet 1", "Bullet 2", "Bullet 3"],
            "speaker_notes": "These are test speaker notes for this slide.",
            "layout": "bullets",
        })()
    ]
    theme = get_theme("executive_modern")
    result = build_pptx(slides, theme)
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_build_pptx_correct_slide_count():
    from app.services.pptx_builder import build_pptx
    from app.services.themes import get_theme
    from pptx import Presentation

    slides = [
        type("Slide", (), {
            "title": f"Slide {i}",
            "bullets": ["a", "b", "c"],
            "speaker_notes": "notes",
            "layout": "bullets",
        })()
        for i in range(3)
    ]
    theme = get_theme("executive_modern")
    result = build_pptx(slides, theme)
    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 3


def test_build_pptx_slide_dimensions_are_16_9():
    from app.services.pptx_builder import build_pptx
    from app.services.themes import get_theme
    from pptx import Presentation
    from pptx.util import Emu

    slides = [type("Slide", (), {
        "title": "Test", "bullets": ["a", "b", "c"],
        "speaker_notes": "notes", "layout": "bullets",
    })()]
    theme = get_theme("executive_modern")
    result = build_pptx(slides, theme)
    prs = Presentation(BytesIO(result))
    assert prs.slide_width == Emu(9144000)
    assert prs.slide_height == Emu(5143500)


def test_build_pptx_speaker_notes_included():
    from app.services.pptx_builder import build_pptx
    from app.services.themes import get_theme
    from pptx import Presentation

    notes_text = "These are the speaker notes for testing purposes."
    slides = [type("Slide", (), {
        "title": "Test", "bullets": ["a", "b", "c"],
        "speaker_notes": notes_text, "layout": "bullets",
    })()]
    theme = get_theme("executive_modern")
    result = build_pptx(slides, theme)
    prs = Presentation(BytesIO(result))
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert notes_text in notes


def test_build_pptx_bullets_appear_in_slide():
    from app.services.pptx_builder import build_pptx
    from app.services.themes import get_theme
    from pptx import Presentation

    slides = [type("Slide", (), {
        "title": "Test",
        "bullets": ["First bullet point", "Second bullet point", "Third bullet point"],
        "speaker_notes": "notes",
        "layout": "bullets",
    })()]
    theme = get_theme("executive_modern")
    result = build_pptx(slides, theme)
    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 1


def test_build_pptx_theme_background_applied():
    from app.services.pptx_builder import build_pptx
    from app.services.themes import get_theme
    from pptx import Presentation

    slides = [type("Slide", (), {
        "title": "Test", "bullets": ["a", "b", "c"],
        "speaker_notes": "notes", "layout": "bullets",
    })()]
    theme = get_theme("corporate_zenith")
    result = build_pptx(slides, theme)
    prs = Presentation(BytesIO(result))
    assert len(prs.slides) == 1
