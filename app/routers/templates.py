import io
import json
import uuid
from datetime import datetime, timezone

import sqlalchemy as sa
from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import get_settings
from app.dependencies.auth import get_current_user, require_admin
from app.dependencies.db import get_db
from app.models.conversion import Conversion
from app.models.slide import Slide
from app.models.template import Template
from app.models.user import User
from app.schemas.template import TemplateCopyResponse, TemplateDetail, TemplateListResponse, TemplateOut
from app.services.audit import log_event
from app.services import gcs as _gcs

router = APIRouter(prefix="/templates", tags=["templates"])

ALLOWED_MIME = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Old binary .ppt format — accepted by some browsers under this MIME type but
# python-pptx cannot parse it; reject early with a clear message.
LEGACY_PPT_MIME = {"application/vnd.ms-powerpoint"}


_DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
# Placeholder indices that represent a slide title
_TITLE_PH_IDX = {0, 13}  # 0=title, 13=center title (common on title slides)


def _extract_shape_texts(shape) -> list[tuple[bool, str]]:
    """Recursively extract (is_title_placeholder, text) from a shape and its children.

    Handles grouped shapes, tables, regular text frames, and SmartArt/complex
    shapes via XML fallback.  Returns a flat list of (is_title, text) tuples.
    """
    results: list[tuple[bool, str]] = []
    try:
        # Grouped shape — recurse into children
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP == 6
            for child in shape.shapes:
                results.extend(_extract_shape_texts(child))
            return results

        # Table — read every cell
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text_frame.text.strip()
                    if text:
                        results.append((False, text))
            return results

        # Regular text frame
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                return results
            is_title = (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx in _TITLE_PH_IDX
            )
            results.append((is_title, text))
            return results

        # Fallback: pull <a:t> text runs directly from XML for SmartArt,
        # charts, and other shapes that don't expose has_text_frame.
        raw_texts = [
            el.text for el in shape.element.iter(f"{{{_DRAWINGML_NS}}}t")
            if el.text and el.text.strip()
        ]
        if raw_texts:
            combined = " ".join(raw_texts).strip()
            if combined:
                results.append((False, combined))
    except Exception:
        pass
    return results


def _parse_pptx_slides(file_bytes: bytes) -> tuple[list[dict], str | None]:
    """Extract slide titles and bullet text from a .pptx file.
    Returns (slides, warning_message). On parse failure returns ([], error_str)."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))

        slides = []
        parse_errors = []

        for i, slide in enumerate(prs.slides):
            title = ""
            bullets: list[str] = []
            notes = ""
            try:
                for shape in slide.shapes:
                    for is_title_ph, text in _extract_shape_texts(shape):
                        if is_title_ph and not title:
                            title = text
                        else:
                            # Split multi-line text blocks into individual bullet lines
                            for line in text.splitlines():
                                line = line.strip()
                                if line and line != title and line not in bullets:
                                    bullets.append(line)

                try:
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text.strip()
                except Exception:
                    pass
            except Exception as e:
                parse_errors.append(str(e))

            slides.append({
                "position": i,
                "layout": "bullets",
                "title": title or f"Slide {i + 1}",
                "bullets": bullets[:8],
                "speaker_notes": notes,
            })

        warning = f"Some slides had parse errors: {'; '.join(parse_errors)}" if parse_errors else None
        return slides, warning
    except Exception as e:
        return [], str(e)


# ── Upload a template (admin → public; user → private) ───────────────────────

@router.post("", status_code=201, response_model=TemplateOut)
async def upload_template(
    name: str = Form(...),
    description: str = Form(""),
    theme: str = Form(""),
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    filename = file.filename or ""
    if file.content_type in LEGACY_PPT_MIME or filename.lower().endswith(".ppt") and not filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=415, detail={"code": "LEGACY_FORMAT", "message": "Old .ppt format is not supported. Please convert to .pptx in PowerPoint or Google Slides and re-upload."})
    if file.content_type not in ALLOWED_MIME and not filename.endswith(".pptx"):
        raise HTTPException(status_code=415, detail={"code": "INVALID_TYPE", "message": "Only .pptx files are accepted"})

    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail={"code": "EMPTY_FILE", "message": "Uploaded file is empty"})
    slides, parse_warning = _parse_pptx_slides(file_bytes)

    # Store the raw file
    template_id = uuid.uuid4()
    gcs_key = f"templates/{template_id}.pptx"

    if _gcs.is_configured():
        settings = get_settings()
        _gcs.upload_bytes(settings.GCS_BUCKET, gcs_key, file_bytes, content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    else:
        local_path = _gcs.LOCAL_UPLOAD_DIR / f"template_{template_id}.pptx"
        _gcs.LOCAL_UPLOAD_DIR.mkdir(exist_ok=True)
        local_path.write_bytes(file_bytes)

    is_admin = current_user.role == "admin"
    template = Template(
        id=template_id,
        name=name,
        description=description or None,
        pptx_key=gcs_key,
        slide_count=len(slides),
        slides_json=slides,
        theme=theme or None,
        is_active=True,
        is_public=is_admin,
        created_by=current_user.id,
    )
    db.add(template)
    await log_event(db, "template.created", actor_id=current_user.id, target_type="template", target_id=template_id,
                    metadata={"name": name, "slides": len(slides), "is_public": is_admin})
    await db.commit()
    await db.refresh(template)
    out = TemplateOut.model_validate(template)
    if parse_warning:
        out.parse_warning = f"Slides could not be extracted ({parse_warning}). Template saved with 0 slides."
    return out


# ── List templates (public + user's own private) ──────────────────────────────

@router.get("", response_model=TemplateListResponse)
async def list_templates(
    page: int = 1,
    page_size: int = 20,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    offset = (page - 1) * page_size
    visibility = sa.or_(Template.is_public == True, Template.created_by == current_user.id)

    total_r = await db.execute(
        sa.select(sa.func.count()).select_from(Template)
        .where(Template.is_active == True, visibility)
    )
    total = total_r.scalar() or 0

    result = await db.execute(
        sa.select(Template)
        .where(Template.is_active == True, visibility)
        .order_by(Template.created_at.desc())
        .offset(offset)
        .limit(page_size)
    )
    items = result.scalars().all()
    return TemplateListResponse(items=[TemplateOut.model_validate(t) for t in items], total=total)


# ── Get template detail ────────────────────────────────────────────────────────

@router.get("/{template_id}", response_model=TemplateDetail)
async def get_template(
    template_id: str,
    db: AsyncSession = Depends(get_db),
    _: User = Depends(get_current_user),
):
    try:
        tid = uuid.UUID(template_id)
    except ValueError:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})

    t = await db.get(Template, tid)
    if not t or not t.is_active:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})
    return TemplateDetail.model_validate(t)


# ── Copy a template → new conversion ─────────────────────────────────────────

@router.post("/{template_id}/copy", response_model=TemplateCopyResponse)
async def copy_template(
    template_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    try:
        tid = uuid.UUID(template_id)
    except ValueError:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})

    t = await db.get(Template, tid)
    if not t or not t.is_active:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})

    slides_data = t.slides_json if isinstance(t.slides_json, list) else []
    slide_count = len(slides_data)

    # Keep original PPTX content as-is — design and content both preserved via build_pptx_from_template on export

    # Normalise any missing fields
    for i, s in enumerate(slides_data):
        if not s.get("title"):
            s["title"] = f"Slide {i + 1}"
        s.setdefault("bullets", [])
        s.setdefault("speaker_notes", "")
        s.setdefault("layout", "bullets")

    # Create conversion record
    conv_id = uuid.uuid4()
    await db.execute(
        sa.text(
            "INSERT INTO conversions (id, user_id, original_filename, status, theme, slide_count, completed_at, source_pptx_key) "
            "VALUES (:id, :uid, :fname, 'done', :theme, :sc, :now, :src_key)"
        ),
        {
            "id": str(conv_id),
            "uid": str(current_user.id),
            "fname": f"{t.name}.pptx",
            "theme": t.theme or "clean_slate",
            "sc": slide_count,
            "now": datetime.now(timezone.utc),
            "src_key": t.pptx_key,
        },
    )

    for slide_data in slides_data:
        await db.execute(
            sa.text(
                "INSERT INTO slides (id, conversion_id, position, layout, title, bullets, speaker_notes, color_scheme, shape_style, is_deleted) "
                "VALUES (:id, :cid, :pos, :layout, :title, CAST(:bullets AS jsonb), :notes, :color_scheme, :shape_style, false)"
            ),
            {
                "id": str(uuid.uuid4()),
                "cid": str(conv_id),
                "pos": slide_data.get("position", 0),
                "layout": slide_data.get("layout", "bullets"),
                "title": slide_data.get("title", ""),
                "bullets": json.dumps(slide_data.get("bullets", [])),
                "notes": slide_data.get("speaker_notes", ""),
                "color_scheme": slide_data.get("color_scheme", "default"),
                "shape_style": slide_data.get("shape_style", "square"),
            },
        )

    await log_event(db, "template.copied", actor_id=current_user.id, target_type="conversion", target_id=conv_id,
                    metadata={"template_id": str(tid), "template_name": t.name})
    await db.commit()

    return TemplateCopyResponse(conversion_id=conv_id, slide_count=slide_count)


# ── Re-parse slides from stored PPTX (admin only) ─────────────────────────────

@router.patch("/{template_id}/reparse", response_model=TemplateOut)
async def reparse_template(
    template_id: str,
    db: AsyncSession = Depends(get_db),
    admin: User = Depends(require_admin),
):
    try:
        tid = uuid.UUID(template_id)
    except ValueError:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})

    t = await db.get(Template, tid)
    if not t or not t.is_active:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})

    file_bytes = _gcs.read_pptx_key_bytes(t.pptx_key)
    if not file_bytes:
        raise HTTPException(status_code=422, detail={"code": "FILE_NOT_FOUND", "message": "Stored PPTX file could not be retrieved"})

    slides, parse_warning = _parse_pptx_slides(file_bytes)

    await db.execute(
        sa.text("UPDATE templates SET slides_json = CAST(:slides AS jsonb), slide_count = :count WHERE id = :id"),
        {"slides": json.dumps(slides), "count": len(slides), "id": str(tid)},
    )
    await log_event(db, "template.reparsed", actor_id=admin.id, target_type="template", target_id=tid,
                    metadata={"name": t.name, "slides": len(slides)})
    await db.commit()
    await db.refresh(t)

    out = TemplateOut.model_validate(t)
    if parse_warning:
        out.parse_warning = f"Slides could not be fully extracted ({parse_warning})."
    return out


# ── Delete a template (admin: any; user: only their own private templates) ─────

@router.delete("/{template_id}", status_code=204)
async def delete_template(
    template_id: str,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    try:
        tid = uuid.UUID(template_id)
    except ValueError:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})

    t = await db.get(Template, tid)
    if not t:
        raise HTTPException(status_code=404, detail={"code": "NOT_FOUND", "message": "Template not found"})

    is_admin = current_user.role == "admin"
    is_owner = t.created_by == current_user.id and not t.is_public
    if not is_admin and not is_owner:
        raise HTTPException(status_code=403, detail={"code": "FORBIDDEN", "message": "You can only delete your own templates"})

    await db.execute(sa.text("UPDATE templates SET is_active = false WHERE id = :id"), {"id": str(tid)})
    await log_event(db, "template.deleted", actor_id=current_user.id, target_type="template", target_id=tid,
                    metadata={"name": t.name})
    await db.commit()
