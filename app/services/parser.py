import io
from dataclasses import dataclass, field
from typing import Any


@dataclass
class Section:
    heading: str
    paragraphs: list[str] = field(default_factory=list)
    bullets: list[str] = field(default_factory=list)
    tables: list[list[list[str]]] = field(default_factory=list)


@dataclass
class ParsedDocument:
    title: str
    sections: list[Section]
    word_count: int
    raw_text: str = ""

    def preview(self, max_words: int = 500) -> str:
        return " ".join(self.raw_text.split()[:max_words])

    def to_prompt_text(self) -> str:
        out = ""
        for s in self.sections:
            out += f"\n## {s.heading}\n"
            for p in s.paragraphs:
                out += f"{p}\n"
            for b in s.bullets:
                out += f"- {b}\n"
        return out

    def to_dict(self) -> dict[str, Any]:
        return {
            "title": self.title,
            "sections": [
                {
                    "heading": s.heading,
                    "paragraphs": s.paragraphs,
                    "bullets": s.bullets,
                    "tables": s.tables,
                }
                for s in self.sections
            ],
            "word_count": self.word_count,
            "raw_text": self.raw_text,
        }

    @classmethod
    def from_dict(cls, data: dict[str, Any]) -> "ParsedDocument":
        sections = [
            Section(
                heading=s["heading"],
                paragraphs=s.get("paragraphs", []),
                bullets=s.get("bullets", []),
                tables=s.get("tables", []),
            )
            for s in data.get("sections", [])
        ]
        return cls(
            title=data.get("title", ""),
            sections=sections,
            word_count=data.get("word_count", 0),
            raw_text=data.get("raw_text", ""),
        )


def _validate(doc: ParsedDocument) -> ParsedDocument:
    if doc.word_count < 30:
        raise ValueError("Document too short")
    if not doc.sections:
        doc.sections = [Section(heading="Document", paragraphs=[doc.raw_text])]
    if len(doc.sections) > 50:
        doc.sections = doc.sections[:50]
    return doc


def parse_pdf(data: bytes) -> ParsedDocument:
    import pdfplumber

    sections: list[Section] = []
    current = Section(heading="Introduction")
    all_text: list[str] = []

    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            all_text.append(text)
            for line in text.split("\n"):
                stripped = line.strip()
                if not stripped:
                    continue
                if len(stripped) < 80 and stripped.isupper():
                    sections.append(current)
                    current = Section(heading=stripped)
                else:
                    current.paragraphs.append(stripped)
            for table in (page.extract_tables() or []):
                if table:
                    current.tables.append(table)

    sections.append(current)
    sections = [s for s in sections if s.paragraphs or s.bullets or s.tables or s.heading != "Introduction"]
    if not sections:
        sections = [current]

    raw_text = "\n".join(all_text)
    word_count = len(raw_text.split())
    title = sections[0].heading if sections else "Document"

    return _validate(ParsedDocument(
        title=title,
        sections=sections,
        word_count=word_count,
        raw_text=raw_text,
    ))


def parse_docx(data: bytes) -> ParsedDocument:
    from docx import Document

    doc = Document(io.BytesIO(data))
    HEADING = {"Heading 1", "Heading 2", "Heading 3", "Title"}
    BULLET = {"List Bullet", "List Bullet 2", "List Paragraph"}

    sections: list[Section] = []
    current = Section(heading="Introduction")
    all_text: list[str] = []

    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        all_text.append(t)
        if para.style.name in HEADING:
            sections.append(current)
            current = Section(heading=t)
        elif para.style.name in BULLET:
            current.bullets.append(t)
        else:
            current.paragraphs.append(t)

    for table in doc.tables:
        rows = [[c.text for c in r.cells] for r in table.rows]
        current.tables.append(rows)

    sections.append(current)
    sections = [s for s in sections if s.paragraphs or s.bullets or s.tables]
    if not sections:
        sections = [current]

    raw_text = "\n".join(all_text)
    word_count = len(raw_text.split())
    title = sections[0].heading if sections else "Document"

    return _validate(ParsedDocument(
        title=title,
        sections=sections,
        word_count=word_count,
        raw_text=raw_text,
    ))


def parse_text(text: str) -> ParsedDocument:
    if not text or not text.strip():
        raise ValueError("Document too short")

    from markdown_it import MarkdownIt

    md = MarkdownIt()
    tokens = md.parse(text)

    sections: list[Section] = []
    current = Section(heading="Document")
    prev_type = ""

    for token in tokens:
        if token.type == "heading_open":
            prev_type = "heading_open"
        elif token.type == "inline" and prev_type == "heading_open":
            sections.append(current)
            current = Section(heading=token.content)
            prev_type = "inline_heading"
        elif token.type == "inline":
            if token.content.strip():
                current.paragraphs.append(token.content)
            prev_type = "inline"
        else:
            if token.type != "heading_open":
                prev_type = token.type

    sections.append(current)
    sections = [s for s in sections if s.paragraphs or s.bullets or s.heading != "Document"]
    if not sections:
        sections = [Section(heading="Document", paragraphs=[text])]

    raw_text = text
    word_count = len(raw_text.split())

    return _validate(ParsedDocument(
        title=sections[0].heading if sections else "Document",
        sections=sections,
        word_count=word_count,
        raw_text=raw_text,
    ))


_DRAWINGML_NS_P = "http://schemas.openxmlformats.org/drawingml/2006/main"
_TITLE_PH_IDX_P = {0, 13}  # 0=title, 13=center title


def _extract_slide_texts(shape) -> list[tuple[bool, str]]:
    """Return (is_title, text) pairs from a shape, recursing into groups."""
    results: list[tuple[bool, str]] = []
    try:
        if shape.shape_type == 6:  # GROUP
            for child in shape.shapes:
                results.extend(_extract_slide_texts(child))
            return results

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    text = cell.text_frame.text.strip()
                    if text:
                        results.append((False, text))
            return results

        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                return results
            is_title = (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx in _TITLE_PH_IDX_P
            )
            results.append((is_title, text))
            return results

        # XML fallback for SmartArt, charts, and other complex shapes
        raw_texts = [
            el.text for el in shape.element.iter(f"{{{_DRAWINGML_NS_P}}}t")
            if el.text and el.text.strip()
        ]
        if raw_texts:
            combined = " ".join(raw_texts).strip()
            if combined:
                results.append((False, combined))
    except Exception:
        pass
    return results


def parse_pptx(data: bytes) -> ParsedDocument:
    import io as _io
    from pptx import Presentation

    prs = Presentation(_io.BytesIO(data))
    sections: list[Section] = []
    all_text: list[str] = []

    try:
        raw_count = len(prs.slides._sldIdLst)
    except Exception:
        raw_count = len(prs.slides)

    for i in range(raw_count):
        try:
            slide = prs.slides[i]
        except Exception:
            continue

        slide_title = ""
        slide_bullets: list[str] = []

        for shape in slide.shapes:
            for is_title_ph, text in _extract_slide_texts(shape):
                if is_title_ph and not slide_title:
                    slide_title = text
                else:
                    for para_line in text.splitlines():
                        line = para_line.strip()
                        if line and line != slide_title:
                            slide_bullets.append(line)

        # Speaker notes as extra context
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_bullets.append(f"[Notes: {notes_text[:300]}]")
        except Exception:
            pass

        heading = slide_title or f"Slide {i + 1}"
        all_text.append(heading)
        all_text.extend(slide_bullets)
        sections.append(Section(heading=heading, bullets=slide_bullets))

    raw_text = "\n".join(all_text)
    word_count = len(raw_text.split())

    if word_count < 10:
        # PPTX has no extractable text (image-only slides)
        raw_text = f"Presentation with {raw_count} slides. No extractable text found."
        word_count = len(raw_text.split())
        sections = [Section(heading="Presentation", paragraphs=[raw_text])]

    title = sections[0].heading if sections else "Presentation"
    return ParsedDocument(title=title, sections=sections, word_count=word_count, raw_text=raw_text)


MIME_PARSERS = {
    "application/pdf": parse_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
    "application/vnd.ms-powerpoint": parse_pptx,
    "text/plain": parse_text,
    "text/markdown": parse_text,
}


def parse_bytes(data: bytes, mime_type: str) -> ParsedDocument:
    parser = MIME_PARSERS.get(mime_type)
    if parser is None:
        raise ValueError(f"Unsupported mime type: {mime_type}")
    if mime_type in ("text/plain", "text/markdown"):
        return parser(data.decode("utf-8", errors="replace"))
    return parser(data)
