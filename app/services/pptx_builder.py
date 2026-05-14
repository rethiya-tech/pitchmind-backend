import io
import logging
from copy import deepcopy
from typing import Any

_log = logging.getLogger(__name__)

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from app.services.themes import Theme

# Slide dimensions — 16:9 widescreen (EMU)
SLIDE_W = 9144000   # 10 in
SLIDE_H = 5143500   # 5.625 in

# Layout constants (EMU) — mirrored exactly in SlidePreview.tsx
HEADER_H   = 685800   # top accent band height  (13.33% of H)
FOOTER_Y   = 5006700  # footer top              (97.34% of H)
FOOTER_H   = 136800   # footer height           ( 2.66% of H)
TITLE_X    = 457200   # title / content left margin
TITLE_Y    = 800000   # title top
TITLE_W    = 8229600  # title width
TITLE_H    = 1143000  # title height
DIVIDER_Y  = 1943400  # thin horizontal divider top
DIVIDER_W  = 1828800  # divider width (2 in)
DIVIDER_H  = 45720    # divider height (0.05 in)
BULLETS_Y  = 2057400  # bullets top
BULLETS_H  = 2743200  # bullets height


# ── helpers ───────────────────────────────────────────────────────────────────

def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex_val(hex_color: str) -> str:
    return hex_color.lstrip("#").upper()


_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
_A = _NS["a"]
_P = _NS["p"]


def _set_background(slide: Any, hex_color: str) -> None:
    """Force-set slide background via direct XML — bypasses bgRef inheritance."""
    color_val = _hex_val(hex_color)
    cSld = slide._element.find(qn("p:cSld"))
    for old_bg in cSld.findall(qn("p:bg")):
        cSld.remove(old_bg)

    bg_elem = etree.fromstring(
        f'<p:bg xmlns:p="{_P}" xmlns:a="{_A}">'
        f'<p:bgPr>'
        f'<a:solidFill><a:srgbClr val="{color_val}"/></a:solidFill>'
        f'<a:effectLst/>'
        f'</p:bgPr>'
        f'</p:bg>'
    )
    spTree = cSld.find(qn("p:spTree"))
    idx = list(cSld).index(spTree) if spTree is not None else 0
    cSld.insert(idx, bg_elem)


def _add_solid_rect(slide: Any, x: int, y: int, w: int, h: int,
                    hex_color: str, shape_id: int) -> None:
    """Insert a coloured rectangle directly via XML for maximum viewer compat."""
    color_val = _hex_val(hex_color)
    sp_xml = (
        f'<p:sp xmlns:p="{_P}" xmlns:a="{_A}">'
        f'  <p:nvSpPr>'
        f'    <p:cNvPr id="{shape_id}" name="rect{shape_id}"/>'
        f'    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'    <p:nvPr/>'
        f'  </p:nvSpPr>'
        f'  <p:spPr>'
        f'    <a:xfrm>'
        f'      <a:off x="{x}" y="{y}"/>'
        f'      <a:ext cx="{w}" cy="{h}"/>'
        f'    </a:xfrm>'
        f'    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'    <a:solidFill><a:srgbClr val="{color_val}"/></a:solidFill>'
        f'    <a:ln><a:noFill/></a:ln>'
        f'  </p:spPr>'
        f'  <p:txBody>'
        f'    <a:bodyPr/><a:lstStyle/><a:p/>'
        f'  </p:txBody>'
        f'</p:sp>'
    )
    spTree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    spTree.append(etree.fromstring(sp_xml))


def _clear_textbox_border(tb: Any) -> None:
    """Remove all borders from a textbox — prevents accent-color bar artifacts."""
    spPr = tb._element.find(qn("p:spPr"))
    if spPr is not None:
        # Remove any existing ln elements (python-pptx may add defaults)
        for ln_elem in list(spPr.findall(qn("a:ln"))):
            spPr.remove(ln_elem)
        # Explicit zero-width no-fill line
        ln = etree.SubElement(spPr, qn("a:ln"), w="0")
        etree.SubElement(ln, qn("a:noFill"))


def _apply_run_style(run: Any, style: dict, default_color: str,
                     default_font: str, default_size: int, default_bold: bool) -> None:
    """Apply a text_styles entry to a run, falling back to provided defaults."""
    color = style.get("color") or default_color
    font_name = style.get("fontFamily") or default_font
    size = style.get("fontSize") or default_size
    weight = style.get("fontWeight")
    bold = (weight >= 600) if weight is not None else default_bold
    italic = bool(style.get("italic"))

    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    run.font.name = font_name


def _add_textbox(slide: Any, text: str, x: int, y: int, w: int, h: int,
                 size_pt: int, bold: bool, hex_color: str,
                 font: str = "Plus Jakarta Sans",
                 style: dict | None = None) -> None:
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    _clear_textbox_border(tb)

    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    run = p.runs[0] if p.runs else p.add_run()
    if style:
        _apply_run_style(run, style, hex_color, font, size_pt, bold)
    else:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = _rgb(hex_color)
        run.font.name = font


def _lighter(hex_color: str, amount: int = 30) -> str:
    h = hex_color.lstrip("#")
    r = min(255, max(0, int(h[0:2], 16) + amount))
    g = min(255, max(0, int(h[2:4], 16) + amount))
    b = min(255, max(0, int(h[4:6], 16) + amount))
    return f"#{r:02X}{g:02X}{b:02X}"


def _add_ellipse(slide: Any, x: int, y: int, w: int, h: int,
                 hex_color: str, shape_id: int) -> None:
    color_val = _hex_val(hex_color)
    sp_xml = (
        f'<p:sp xmlns:p="{_P}" xmlns:a="{_A}">'
        f'  <p:nvSpPr>'
        f'    <p:cNvPr id="{shape_id}" name="ellipse{shape_id}"/>'
        f'    <p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        f'    <p:nvPr/>'
        f'  </p:nvSpPr>'
        f'  <p:spPr>'
        f'    <a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{w}" cy="{h}"/></a:xfrm>'
        f'    <a:prstGeom prst="ellipse"><a:avLst/></a:prstGeom>'
        f'    <a:solidFill><a:srgbClr val="{color_val}"/></a:solidFill>'
        f'    <a:ln><a:noFill/></a:ln>'
        f'  </p:spPr>'
        f'  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        f'</p:sp>'
    )
    spTree = slide._element.find(qn("p:cSld")).find(qn("p:spTree"))
    spTree.append(etree.fromstring(sp_xml))


def _add_bullet_column(sl: Any, bullets: list[str], x: int, y: int, w: int, h: int,
                       theme: Any, font: str = "Plus Jakarta Sans",
                       bullet_styles: dict | None = None,
                       bullet_offset: int = 0) -> None:
    """Add a column of bullet points as a textbox.

    bullet_styles: dict keyed by str(absolute_bullet_index) → SlideTextStyle dict
    bullet_offset: index of first bullet in this column (for two_column layouts)
    """
    cb = sl.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    _clear_textbox_border(cb)
    ctf = cb.text_frame
    ctf.word_wrap = True
    styles = bullet_styles or {}
    for j, bullet in enumerate(bullets):
        abs_idx = j + bullet_offset
        cp = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
        dot = cp.add_run()
        dot.text = "● "
        dot.font.size = Pt(9)
        dot.font.color.rgb = _rgb(theme.accent)
        dot.font.name = font
        txt = cp.add_run()
        txt.text = bullet
        bstyle = styles.get(str(abs_idx))
        if bstyle:
            _apply_run_style(txt, bstyle, theme.text, font, 13, False)
        else:
            txt.font.size = Pt(13)
            txt.font.color.rgb = _rgb(theme.text)
            txt.font.name = font


# ── layout renderers ──────────────────────────────────────────────────────────

def _render_common_header(sl: Any, theme: Any, slide_idx: int, sid: int) -> None:
    """Shared header band + slide number chip."""
    num_w = 457200
    _add_solid_rect(sl, 0, 0, SLIDE_W, HEADER_H, theme.accent, sid)
    _add_solid_rect(sl, SLIDE_W - num_w, 0, num_w, HEADER_H, theme.accent, sid + 1)
    _add_textbox(sl, str(slide_idx + 1), SLIDE_W - num_w, 0, num_w, HEADER_H,
                 size_pt=10, bold=True, hex_color="#FFFFFF")


def _render_common_footer(sl: Any, theme: Any, sid: int) -> None:
    _add_solid_rect(sl, 0, FOOTER_Y, SLIDE_W, FOOTER_H, theme.accent, sid)
    _add_textbox(sl, "PitchMind", TITLE_X, FOOTER_Y, 2000000, FOOTER_H,
                 size_pt=8, bold=False, hex_color="#FFFFFF")


def _render_hero(sl: Any, slide: Any, theme: Any, sid: int) -> None:
    ts = getattr(slide, 'text_styles', None) or {}
    _add_textbox(sl, slide.title or "", TITLE_X, 1300000, TITLE_W, 1600000,
                 size_pt=38, bold=True, hex_color=theme.text,
                 style=ts.get('title'))
    _add_solid_rect(sl, TITLE_X, 3000000, DIVIDER_W * 2, DIVIDER_H * 2, theme.accent, sid)
    subtitle = slide.bullets[0] if slide.bullets else ""
    if subtitle:
        bullet_styles = ts.get('bullets') or {}
        _add_textbox(sl, subtitle, TITLE_X, 3200000, TITLE_W, 700000,
                     size_pt=18, bold=False, hex_color=theme.text,
                     style=bullet_styles.get('0'))
    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


# No-header layout constants (EMU)
_NH_TITLE_Y   = 380000   # ~7% of H
_NH_TITLE_H   = 1200000
_NH_DIVIDER_Y = 1650000  # ~32%
_NH_BULLETS_Y = 1750000  # ~34%
_NH_BULLETS_H = 3200000  # to ~96%


def _render_bullets(sl: Any, slide: Any, theme: Any, sid: int, slide_idx: int) -> None:
    ts = getattr(slide, 'text_styles', None) or {}
    _add_textbox(sl, slide.title or "", TITLE_X, _NH_TITLE_Y, TITLE_W, _NH_TITLE_H,
                 size_pt=28, bold=True, hex_color=theme.text,
                 style=ts.get('title'))
    _add_solid_rect(sl, TITLE_X, _NH_DIVIDER_Y, DIVIDER_W, DIVIDER_H, theme.accent, sid)
    bullets = slide.bullets if slide.bullets else []
    if bullets:
        _add_bullet_column(sl, bullets, TITLE_X, _NH_BULLETS_Y, TITLE_W, _NH_BULLETS_H,
                           theme, bullet_styles=ts.get('bullets'))
    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


def _render_two_column(sl: Any, slide: Any, theme: Any, sid: int, slide_idx: int) -> None:
    ts = getattr(slide, 'text_styles', None) or {}
    _add_textbox(sl, slide.title or "", TITLE_X, _NH_TITLE_Y, TITLE_W, _NH_TITLE_H,
                 size_pt=24, bold=True, hex_color=theme.text,
                 style=ts.get('title'))
    _add_solid_rect(sl, TITLE_X, _NH_DIVIDER_Y, DIVIDER_W, DIVIDER_H, theme.accent, sid)

    bullets = slide.bullets if slide.bullets else []
    mid = (len(bullets) + 1) // 2
    left_bullets = bullets[:mid]
    right_bullets = bullets[mid:]

    gap = 228600
    col_w = (TITLE_W - gap) // 2
    right_x = TITLE_X + col_w + gap

    col_hdr_h = 228600
    col_content_y = _NH_BULLETS_Y + col_hdr_h + 60000
    col_content_h = _NH_BULLETS_H - col_hdr_h - 60000

    _add_solid_rect(sl, TITLE_X, _NH_BULLETS_Y, col_w, col_hdr_h, theme.accent, sid + 1)
    _add_textbox(sl, "Key Points", TITLE_X + 80000, _NH_BULLETS_Y, col_w - 160000, col_hdr_h,
                 size_pt=9, bold=True, hex_color="#FFFFFF")

    _add_solid_rect(sl, right_x, _NH_BULLETS_Y, col_w, col_hdr_h, _lighter(theme.accent, 20), sid + 2)
    _add_textbox(sl, "Details", right_x + 80000, _NH_BULLETS_Y, col_w - 160000, col_hdr_h,
                 size_pt=9, bold=True, hex_color="#FFFFFF")

    bullet_styles = ts.get('bullets')
    if left_bullets:
        _add_bullet_column(sl, left_bullets, TITLE_X, col_content_y, col_w, col_content_h,
                           theme, bullet_styles=bullet_styles, bullet_offset=0)
    if right_bullets:
        _add_bullet_column(sl, right_bullets, right_x, col_content_y, col_w, col_content_h,
                           theme, bullet_styles=bullet_styles, bullet_offset=mid)

    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


def _render_data_table(sl: Any, slide: Any, theme: Any, sid: int, slide_idx: int) -> None:
    _add_textbox(sl, slide.title or "", TITLE_X, _NH_TITLE_Y, TITLE_W, _NH_TITLE_H,
                 size_pt=24, bold=True, hex_color=theme.text)

    bullets = slide.bullets if slide.bullets else []
    table_y = _NH_BULLETS_Y
    row_h = 350000
    gap = 22000
    label_w = int(TITLE_W * 0.38)
    value_w = TITLE_W - label_w
    padding = 70000
    value_bg = _lighter(theme.bg, 35)

    for i, bullet in enumerate(bullets):
        label, value = bullet.split(': ', 1) if ': ' in bullet else (bullet, "")
        y = table_y + i * (row_h + gap)
        label_bg = theme.accent if i % 2 == 0 else _lighter(theme.accent, 25)
        base = sid + i * 4

        _add_solid_rect(sl, TITLE_X, y, label_w, row_h, label_bg, base)
        _add_textbox(sl, label, TITLE_X + padding, y + padding // 3, label_w - padding * 2, row_h - padding // 2,
                     size_pt=11, bold=True, hex_color="#FFFFFF")
        _add_solid_rect(sl, TITLE_X + label_w, y, value_w, row_h, value_bg, base + 1)
        _add_textbox(sl, value or "—", TITLE_X + label_w + padding, y + padding // 3,
                     value_w - padding * 2, row_h - padding // 2,
                     size_pt=11, bold=False, hex_color=theme.text)

    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


def _render_timeline(sl: Any, slide: Any, theme: Any, sid: int, slide_idx: int) -> None:
    _add_textbox(sl, slide.title or "", TITLE_X, _NH_TITLE_Y, TITLE_W, _NH_TITLE_H,
                 size_pt=24, bold=True, hex_color=theme.text)

    bullets = slide.bullets if slide.bullets else []
    n = min(len(bullets), 6)
    if n == 0:
        if slide.speaker_notes:
            sl.notes_slide.notes_text_frame.text = slide.speaker_notes
        return

    node_r = 130000
    node_d = node_r * 2
    line_cx = TITLE_X + node_r
    content_x = TITLE_X + node_d + 150000
    content_w = TITLE_W - node_d - 150000
    row_h = _NH_BULLETS_H // n

    # Vertical connecting line
    _add_solid_rect(sl, line_cx - 18000, _NH_BULLETS_Y + node_r,
                    36000, row_h * n - node_d, theme.accent, sid)

    for i, bullet in enumerate(bullets[:n]):
        parts = bullet.split(': ', 1)
        label = parts[0] if len(parts) == 2 else f"Phase {i + 1}"
        desc = parts[1] if len(parts) == 2 else bullet
        y = _NH_BULLETS_Y + i * row_h
        base = sid + 2 + i * 5
        _add_ellipse(sl, line_cx - node_r, y, node_d, node_d, theme.accent, base)
        label_h = 380000
        _add_textbox(sl, label, content_x, y, content_w, label_h,
                     size_pt=11, bold=True, hex_color=theme.accent)
        if desc:
            _add_textbox(sl, desc, content_x, y + label_h, content_w, row_h - label_h - 30000,
                         size_pt=10, bold=False, hex_color=theme.text)

    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


def _render_big_stat(sl: Any, slide: Any, theme: Any, sid: int, slide_idx: int) -> None:
    _add_textbox(sl, slide.title or "", TITLE_X, _NH_TITLE_Y, TITLE_W, _NH_TITLE_H,
                 size_pt=24, bold=True, hex_color=theme.text)

    bullets = slide.bullets if slide.bullets else []
    n = min(len(bullets), 3)
    if n == 0:
        if slide.speaker_notes:
            sl.notes_slide.notes_text_frame.text = slide.speaker_notes
        return

    gap = 228600
    card_w = (TITLE_W - gap * (n - 1)) // n
    card_h = 2400000
    card_y = _NH_BULLETS_Y + 200000
    accent_bar_w = 80000
    padding = 150000

    for i, bullet in enumerate(bullets[:n]):
        parts = bullet.split(': ', 1)
        label = parts[0] if len(parts) == 2 else "Metric"
        value = parts[1] if len(parts) == 2 else bullet
        x = TITLE_X + i * (card_w + gap)
        base = sid + 2 + i * 5
        _add_solid_rect(sl, x, card_y, card_w, card_h, _lighter(theme.bg, 30), base)
        _add_solid_rect(sl, x, card_y, accent_bar_w, card_h, theme.accent, base + 1)
        _add_textbox(sl, value, x + accent_bar_w + padding, card_y + 200000,
                     card_w - accent_bar_w - padding * 2, 1200000,
                     size_pt=52, bold=True, hex_color=theme.accent)
        _add_textbox(sl, label, x + accent_bar_w + padding, card_y + 1500000,
                     card_w - accent_bar_w - padding * 2, 600000,
                     size_pt=14, bold=False, hex_color=theme.text)

    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


def _render_process(sl: Any, slide: Any, theme: Any, sid: int, slide_idx: int) -> None:
    _add_textbox(sl, slide.title or "", TITLE_X, _NH_TITLE_Y, TITLE_W, _NH_TITLE_H,
                 size_pt=24, bold=True, hex_color=theme.text)
    _add_solid_rect(sl, TITLE_X, _NH_DIVIDER_Y, DIVIDER_W, DIVIDER_H, theme.accent, sid)

    bullets = slide.bullets if slide.bullets else []
    n = min(len(bullets), 5)
    if n == 0:
        if slide.speaker_notes:
            sl.notes_slide.notes_text_frame.text = slide.speaker_notes
        return

    arrow_w = 150000
    gap = 60000
    step_span = arrow_w + gap * 2
    box_w = (TITLE_W - step_span * (n - 1)) // n
    box_h = 2000000
    box_y = _NH_BULLETS_Y + 200000
    num_h = 500000
    dark_accent = _lighter(theme.accent, -20)

    for i, bullet in enumerate(bullets[:n]):
        x = TITLE_X + i * (box_w + step_span)
        base = sid + 2 + i * 5
        box_bg = theme.accent if i % 2 == 0 else _lighter(theme.accent, 30)
        _add_solid_rect(sl, x, box_y, box_w, box_h, box_bg, base)
        _add_solid_rect(sl, x, box_y, box_w, num_h, dark_accent, base + 1)
        _add_textbox(sl, str(i + 1), x, box_y, box_w, num_h,
                     size_pt=18, bold=True, hex_color="#FFFFFF")
        _add_textbox(sl, bullet, x + 60000, box_y + num_h + 60000,
                     box_w - 120000, box_h - num_h - 120000,
                     size_pt=11, bold=False, hex_color="#FFFFFF")
        if i < n - 1:
            arrow_x = x + box_w + gap
            arrow_y = box_y + box_h // 2 - 100000
            _add_textbox(sl, "→", arrow_x, arrow_y, arrow_w, 300000,
                         size_pt=16, bold=True, hex_color=theme.accent)

    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


def _render_quote(sl: Any, slide: Any, theme: Any, sid: int, slide_idx: int) -> None:
    _add_textbox(sl, slide.title or "", TITLE_X, _NH_TITLE_Y, TITLE_W, 800000,
                 size_pt=18, bold=False, hex_color=theme.text)

    bullets = slide.bullets if slide.bullets else []
    quote_text = bullets[0] if bullets else ""
    attribution = bullets[1] if len(bullets) > 1 else ""

    # Large opening quotation mark
    _add_textbox(sl, "“", TITLE_X, 1200000, 800000, 1200000,
                 size_pt=80, bold=True, hex_color=theme.accent)

    if quote_text:
        _add_textbox(sl, quote_text, TITLE_X + 600000, 1400000,
                     TITLE_W - 600000, 2200000,
                     size_pt=22, bold=False, hex_color=theme.text)

    _add_solid_rect(sl, TITLE_X, 3800000, DIVIDER_W * 2, DIVIDER_H * 2, theme.accent, sid)

    if attribution:
        _add_textbox(sl, f"— {attribution}", SLIDE_W - 4000000, 3900000, 3600000, 600000,
                     size_pt=13, bold=False, hex_color=theme.text)

    if slide.speaker_notes:
        sl.notes_slide.notes_text_frame.text = slide.speaker_notes


# ── main builder ──────────────────────────────────────────────────────────────

def _purge_spTree(element: Any, label: str) -> None:
    """Remove ALL non-placeholder shapes from a cSld element's spTree."""
    spTree = element.find(qn("p:cSld")).find(qn("p:spTree"))
    all_children = list(spTree)
    _log.warning("PPTX DEBUG %s spTree children: %s",
                 label, [c.tag.split("}")[-1] for c in all_children])
    for sp in list(spTree.findall(qn("p:sp"))):
        is_ph = False
        try:
            nvPr = sp.find(qn("p:nvSpPr")).find(qn("p:nvPr"))
            if nvPr is not None and nvPr.find(qn("p:ph")) is not None:
                is_ph = True
        except AttributeError:
            pass
        _log.warning("PPTX DEBUG %s sp is_ph=%s", label, is_ph)
        if not is_ph:
            spTree.remove(sp)
    for tag in (qn("p:pic"), qn("p:graphicFrame"), qn("p:cxnSp"), qn("p:grpSp")):
        for elem in list(spTree.findall(tag)):
            _log.warning("PPTX DEBUG %s removing %s", label, tag.split("}")[-1])
            spTree.remove(elem)


def _blank_layout(prs: Any) -> Any:
    """Return a clean Blank slide layout with master shapes also stripped."""
    _log.warning("PPTX DEBUG layouts: %s", [l.name for l in prs.slide_layouts])
    layout = None
    for l in prs.slide_layouts:
        if l.name.strip().lower() == "blank":
            layout = l
            break
    if layout is None:
        _log.warning("PPTX DEBUG no 'Blank' layout found, using index 6")
        layout = prs.slide_layouts[6]
    _log.warning("PPTX DEBUG using layout: '%s'", layout.name)
    _purge_spTree(layout._element, f"layout[{layout.name}]")
    _purge_spTree(layout.slide_master._element, "master")

    # Zero out the accent color in the master theme so viewers can't inherit
    # an accent-colored border on shapes that don't explicitly set their line.
    try:
        master = layout.slide_master
        theme_elem = master._element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )
        if theme_elem is None:
            theme_elem = master._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}fmtScheme"
            )
        # Override every ln (line) style in the master's format scheme to noFill
        for ln in master._element.iter(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}ln"
        ):
            for child in list(ln):
                ln.remove(child)
            ln.set("w", "0")
            etree.SubElement(
                ln, "{http://schemas.openxmlformats.org/drawingml/2006/main}noFill"
            )
    except Exception:
        pass  # Non-critical — explicit noFill on each shape is the primary fix

    return layout


def build_pptx(slides: list[Any], theme: Theme) -> bytes:
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    blank = _blank_layout(prs)  # layout cleaned once here

    for slide_idx, slide in enumerate(slides):
        sl = prs.slides.add_slide(blank)
        sl._element.set("showMasterSp", "0")

        # Fix empty grpSpPr — third-party viewers require an explicit xfrm on
        # the spTree's group so they know the child coordinate space matches the
        # slide dimensions.  Without it, viewers like Jumpshare/Doconut render
        # shapes in unexpected positions.
        cSld = sl._element.find(qn("p:cSld"))
        spTree = cSld.find(qn("p:spTree"))
        grpSpPr = spTree.find(qn("p:grpSpPr"))
        if grpSpPr is not None and grpSpPr.find(qn("a:xfrm")) is None:
            xfrm = etree.SubElement(grpSpPr, qn("a:xfrm"))
            etree.SubElement(xfrm, qn("a:off"), x="0", y="0")
            etree.SubElement(xfrm, qn("a:ext"), cx=str(SLIDE_W), cy=str(SLIDE_H))
            etree.SubElement(xfrm, qn("a:chOff"), x="0", y="0")
            etree.SubElement(xfrm, qn("a:chExt"), cx=str(SLIDE_W), cy=str(SLIDE_H))

        # 1. Background — image if available, else solid color
        bg_path = theme.bg_image_path()
        if bg_path:
            sl.shapes.add_picture(bg_path, Emu(0), Emu(0), Emu(SLIDE_W), Emu(SLIDE_H))
        else:
            _set_background(sl, theme.bg)

        sid = slide_idx * 20 + 10
        layout = getattr(slide, 'layout', 'bullets') or 'bullets'

        if layout == 'hero':
            _render_hero(sl, slide, theme, sid)
        elif layout == 'two_column':
            _render_two_column(sl, slide, theme, sid, slide_idx)
        elif layout == 'data_table':
            _render_data_table(sl, slide, theme, sid, slide_idx)
        elif layout == 'timeline':
            _render_timeline(sl, slide, theme, sid, slide_idx)
        elif layout == 'big_stat':
            _render_big_stat(sl, slide, theme, sid, slide_idx)
        elif layout == 'process':
            _render_process(sl, slide, theme, sid, slide_idx)
        elif layout == 'quote':
            _render_quote(sl, slide, theme, sid, slide_idx)
        else:
            _render_bullets(sl, slide, theme, sid, slide_idx)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Template-preserving builder ───────────────────────────────────────────────

_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _delete_slide(prs: Any, idx: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    rId = sld_id_lst[idx].get(f"{{{_R_NS}}}id")
    if rId:
        prs.part.drop_rel(rId)
    del sld_id_lst[idx]


def _clone_slide(prs: Any, source_idx: int) -> None:
    """Append a deep copy of the slide at source_idx to the end of the deck."""
    src_slide = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(src_slide.slide_layout)
    src_cSld = src_slide._element.find(qn("p:cSld"))
    tgt_cSld = new_slide._element.find(qn("p:cSld"))
    if src_cSld is not None and tgt_cSld is not None:
        parent = tgt_cSld.getparent()
        idx_in_parent = list(parent).index(tgt_cSld)
        parent.remove(tgt_cSld)
        parent.insert(idx_in_parent, deepcopy(src_cSld))


def _set_tf_text(tf: Any, lines: list[str]) -> None:
    """Replace a text frame's content with lines, preserving first-run formatting."""
    first_runs = tf.paragraphs[0].runs if tf.paragraphs else []
    sample_run_xml = deepcopy(first_runs[0]._r) if first_runs else None
    txBody = tf._txBody
    for p_elem in list(txBody.findall(qn("a:p"))):
        txBody.remove(p_elem)
    for line in (lines if lines else [""]):
        p_elem = etree.SubElement(txBody, qn("a:p"))
        r_elem = deepcopy(sample_run_xml) if sample_run_xml is not None else etree.SubElement(p_elem, qn("a:r"))
        t_elem = r_elem.find(qn("a:t"))
        if t_elem is None:
            t_elem = etree.SubElement(r_elem, qn("a:t"))
        t_elem.text = line
        p_elem.append(r_elem)


def _update_slide_placeholders(slide: Any, title: str, bullets: list[str], notes: str) -> None:
    """Update title and body placeholders in a slide; leave all visuals intact.

    Falls back to updating the two largest textboxes (by area) when the slide
    has no title/body placeholder shapes — common in hand-crafted PPTX files.
    """
    title_updated = False
    body_updated = False

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if not shape.is_placeholder:
            continue
        ph_idx = shape.placeholder_format.idx
        tf = shape.text_frame

        if ph_idx == 0:
            p0 = tf.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = title
                for run in p0.runs[1:]:
                    run.text = ""
            else:
                p0.text = title
            for extra_p in list(tf.paragraphs[1:]):
                extra_p._p.getparent().remove(extra_p._p)
            title_updated = True

        elif ph_idx in (1, 2, 13, 14, 15):
            _set_tf_text(tf, bullets)
            body_updated = True

    # Fallback: no placeholders found — update textboxes sorted by vertical position
    if not title_updated or not body_updated:
        textboxes = sorted(
            [s for s in slide.shapes if s.has_text_frame and not s.is_placeholder],
            key=lambda s: (s.top, s.left),
        )
        if not title_updated and textboxes:
            tb = textboxes[0]
            p0 = tb.text_frame.paragraphs[0]
            if p0.runs:
                p0.runs[0].text = title
                for run in p0.runs[1:]:
                    run.text = ""
            else:
                p0.text = title
            for extra_p in list(tb.text_frame.paragraphs[1:]):
                extra_p._p.getparent().remove(extra_p._p)
        if not body_updated and len(textboxes) > 1:
            _set_tf_text(textboxes[1].text_frame, bullets)

    if notes:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass


def build_pptx_from_template(user_slides: list[Any], template_bytes: bytes) -> bytes:
    """Export using the original PPTX as the visual base, updating only text."""
    prs = Presentation(io.BytesIO(template_bytes))

    template_count = len(prs.slides._sldIdLst)
    user_count = len(user_slides)

    # Extend: clone last template slide to fill missing positions
    while len(prs.slides._sldIdLst) < user_count:
        _clone_slide(prs, template_count - 1)

    # Shrink: remove slides from the end
    while len(prs.slides._sldIdLst) > user_count:
        _delete_slide(prs, len(prs.slides._sldIdLst) - 1)

    # Update text in each slide
    for i, user_slide in enumerate(user_slides):
        slide = prs.slides[i]
        title = getattr(user_slide, "title", "") or ""
        bullets = list(getattr(user_slide, "bullets", []) or [])
        notes = getattr(user_slide, "speaker_notes", "") or ""
        _update_slide_placeholders(slide, title, bullets, notes)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
